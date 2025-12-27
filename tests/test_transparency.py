"""Tests for transparency handling in image processing."""

import io
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pp.core.image_processor import apply_color_transform


FIXTURES_DIR = Path(__file__).parent / "fixtures"
HAGAR_PPTX = FIXTURES_DIR / "hagar-presentatie.pptx"


class TestRGBWithTransparencyColor:
    """Tests for RGB images that have transparency defined via info['transparency'].
    
    These are PNG images that use a specific RGB color to indicate transparency,
    rather than having an explicit alpha channel (RGBA).
    """

    @pytest.fixture
    def rgb_image_with_transparency(self) -> Image.Image:
        """Create an RGB image with transparency color set to white.
        
        This mimics the format found in hagar-presentatie.pptx where:
        - Mode is RGB (not RGBA)
        - White pixels (255, 255, 255) should be transparent
        - Transparency is indicated via img.info['transparency']
        """
        # Create a 100x100 image with black text on white background
        img = Image.new("RGB", (100, 100), (255, 255, 255))  # White background
        
        # Draw some black pixels (simulating text)
        pixels = img.load()
        for x in range(20, 80):
            for y in range(40, 60):
                pixels[x, y] = (0, 0, 0)  # Black text
        
        # Set transparency color to white (this is how PNG transparency works)
        img.info['transparency'] = (255, 255, 255)
        
        return img

    def test_transparency_color_preserved_after_inversion(self, rgb_image_with_transparency: Image.Image):
        """Test that RGB images with transparency color are converted to RGBA with proper alpha.
        
        When inverting an RGB image with transparency:
        1. White (transparent) pixels should become transparent in the output
        2. Black (opaque) pixels should become the target light color and remain opaque
        3. The output should be RGBA to properly represent transparency
        """
        img = rgb_image_with_transparency
        
        # Verify preconditions
        assert img.mode == "RGB"
        assert "transparency" in img.info
        assert img.info["transparency"] == (255, 255, 255)
        
        # Apply standard black/white inversion
        result = apply_color_transform(
            img,
            target_dark=(0, 0, 0),      # Black background
            target_light=(255, 255, 255) # White foreground
        )
        
        # The result should be RGBA to preserve transparency
        assert result.mode == "RGBA", "Result should be RGBA to preserve transparency"
        
        # Check alpha channel
        alpha = result.split()[3]
        list(alpha.getdata())
        
        # Originally white pixels (background) should now be transparent (alpha=0)
        # Originally black pixels (text) should be opaque (alpha=255)
        
        # Check a pixel that was white (should be transparent)
        # Position (0, 0) was white background
        assert result.getpixel((0, 0))[3] == 0, "Originally white pixels should be transparent"
        
        # Check a pixel that was black (should be opaque)  
        # Position (50, 50) was black text
        assert result.getpixel((50, 50))[3] == 255, "Originally black pixels should be opaque"

    def test_inverted_colors_correct_with_transparency(self, rgb_image_with_transparency: Image.Image):
        """Test that colors are correctly inverted while preserving transparency.
        
        Black text should become white, transparent areas should remain transparent.
        """
        img = rgb_image_with_transparency
        
        result = apply_color_transform(
            img,
            target_dark=(0, 0, 0),
            target_light=(255, 255, 255)
        )
        
        # Black pixels (text at 50, 50) should become white
        pixel = result.getpixel((50, 50))
        assert pixel[:3] == (255, 255, 255), f"Black should invert to white, got {pixel[:3]}"
        assert pixel[3] == 255, "Text pixels should be fully opaque"

    def test_custom_colors_with_transparency(self, rgb_image_with_transparency: Image.Image):
        """Test custom color inversion preserves transparency."""
        img = rgb_image_with_transparency
        
        # Invert to yellow text on navy background
        result = apply_color_transform(
            img,
            target_dark=(0, 0, 128),     # Navy (for originally light/transparent areas)
            target_light=(255, 255, 0)   # Yellow (for originally dark areas)
        )
        
        assert result.mode == "RGBA"
        
        # Originally black text should become yellow and opaque
        text_pixel = result.getpixel((50, 50))
        assert text_pixel[:3] == (255, 255, 0), f"Text should be yellow, got {text_pixel[:3]}"
        assert text_pixel[3] == 255, "Text should be opaque"
        
        # Originally white background should be transparent
        bg_pixel = result.getpixel((0, 0))
        assert bg_pixel[3] == 0, "Background should be transparent"


class TestHagarPresentatieImages:
    """Tests using the actual hagar-presentatie.pptx fixture."""

    @pytest.mark.skipif(not HAGAR_PPTX.exists(), reason="hagar-presentatie.pptx not found")
    def test_extract_image_has_transparency_info(self):
        """Verify the fixture images have the expected transparency format."""
        prs = Presentation(str(HAGAR_PPTX))
        
        # Check the second slide (index 1) which shows "slacht,"
        slide = prs.slides[1]
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                with Image.open(io.BytesIO(shape.image.blob)) as img:
                    assert img.mode == "RGB", "Expected RGB mode"
                    assert "transparency" in img.info, "Expected transparency info"
                    assert img.info["transparency"] == (255, 255, 255), "Expected white transparency"

    @pytest.mark.skipif(not HAGAR_PPTX.exists(), reason="hagar-presentatie.pptx not found")
    def test_invert_hagar_image_preserves_transparency(self):
        """Test that inverting images from hagar-presentatie.pptx preserves transparency.
        
        The images have black text on transparent (white) background.
        After inversion, we should get white text with transparency preserved.
        """
        prs = Presentation(str(HAGAR_PPTX))
        slide = prs.slides[1]
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                with Image.open(io.BytesIO(shape.image.blob)) as img:
                    result = apply_color_transform(
                        img,
                        target_dark=(0, 0, 0),
                        target_light=(255, 255, 255)
                    )
                    
                    # Result must be RGBA to preserve transparency
                    assert result.mode == "RGBA", "Result should be RGBA"
                    
                    # Check that we have both transparent and opaque pixels
                    alpha = result.split()[3]
                    alpha_extrema = alpha.getextrema()
                    
                    assert alpha_extrema[0] == 0, "Should have fully transparent pixels"
                    assert alpha_extrema[1] == 255, "Should have fully opaque pixels"
                    
                    # Count transparent pixels - should be significant (background)
                    alpha_data = list(alpha.getdata())
                    transparent_count = sum(1 for a in alpha_data if a == 0)
                    total_pixels = len(alpha_data)
                    
                    # The background is white/transparent, which is most of the image
                    assert transparent_count > total_pixels * 0.5, \
                        f"Expected >50% transparent pixels, got {transparent_count/total_pixels*100:.1f}%"
