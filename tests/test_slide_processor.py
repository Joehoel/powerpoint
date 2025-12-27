"""Tests for slide processing."""

import io
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor

from pp.core.slide_processor import process_slide, process_slide_safe
from pp.models.config import InversionConfig


FIXTURES_DIR = Path(__file__).parent / "fixtures"


class TestProcessSlide:
    """Tests for the process_slide function."""

    def test_process_slide_background(self, sample_presentation: Presentation, default_config: InversionConfig):
        """Test that slide background is set correctly."""
        slide = sample_presentation.slides[0]
        
        process_slide(slide, default_config)
        
        # Background should be black
        fill = slide.background.fill
        assert fill.fore_color.rgb == RGBColor(0, 0, 0)

    def test_process_slide_text_color(self, sample_presentation: Presentation, default_config: InversionConfig):
        """Test that text color is changed to foreground color."""
        slide = sample_presentation.slides[0]
        
        process_slide(slide, default_config)
        
        # Find shapes with text and check color
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        assert run.font.color.rgb == RGBColor(255, 255, 255)

    def test_process_slide_custom_colors(self, sample_presentation: Presentation, custom_config: InversionConfig):
        """Test processing with custom colors."""
        slide = sample_presentation.slides[0]
        
        process_slide(slide, custom_config)
        
        # Background should be navy
        fill = slide.background.fill
        assert fill.fore_color.rgb == RGBColor(0, 0, 128)
        
        # Text should be yellow
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        assert run.font.color.rgb == RGBColor(255, 255, 0)

    def test_process_slide_with_image(self, sample_presentation: Presentation, default_config: InversionConfig):
        """Test processing slide with images."""
        # Slide 2 has an image
        slide = sample_presentation.slides[1]
        
        process_slide(slide, default_config)
        
        # Should complete without errors (warnings may exist)
        # Background should still be set
        fill = slide.background.fill
        assert fill.fore_color.rgb == RGBColor(0, 0, 0)

    def test_process_slide_no_image_inversion(self, sample_presentation: Presentation):
        """Test processing with image inversion disabled."""
        config = InversionConfig(invert_images=False)
        slide = sample_presentation.slides[1]
        
        # Count shapes before
        shape_count_before = len(list(slide.shapes))
        
        process_slide(slide, config)
        
        # Shape count should be same (image not replaced)
        shape_count_after = len(list(slide.shapes))
        assert shape_count_after == shape_count_before


class TestProcessSlideSafe:
    """Tests for the process_slide_safe function."""

    def test_returns_success_on_valid_slide(self, sample_presentation: Presentation, default_config: InversionConfig):
        """Test that valid slides return success."""
        slide = sample_presentation.slides[0]
        
        success, warnings = process_slide_safe(slide, default_config)
        
        assert success is True

    def test_collects_warnings(self, sample_presentation: Presentation, default_config: InversionConfig):
        """Test that warnings are collected."""
        slide = sample_presentation.slides[1]  # Slide with image
        
        success, warnings = process_slide_safe(slide, default_config)
        
        # Should succeed even if there are warnings
        assert success is True
        assert isinstance(warnings, list)


class TestWithRealFixtures:
    """Tests using real PPTX fixture files."""

    @pytest.mark.skipif(not FIXTURES_DIR.exists(), reason="Fixtures directory not found")
    def test_process_real_pptx(self, default_config: InversionConfig):
        """Test processing a real PPTX file."""
        fixture_files = list(FIXTURES_DIR.glob("*.pptx"))
        if not fixture_files:
            pytest.skip("No PPTX fixtures found")
        
        pptx_path = fixture_files[0]
        prs = Presentation(str(pptx_path))
        
        all_warnings = []
        for slide in prs.slides:
            warnings = process_slide(slide, default_config)
            all_warnings.extend(warnings)
        
        # Should complete without crashing
        # Save to verify it's still valid
        output = io.BytesIO()
        prs.save(output)
        
        # Should be able to re-open
        output.seek(0)
        reopened = Presentation(output)
        assert len(reopened.slides) == len(prs.slides)

    @pytest.mark.skipif(not FIXTURES_DIR.exists(), reason="Fixtures directory not found")
    def test_process_multiple_fixtures(self, default_config: InversionConfig):
        """Test processing multiple PPTX files."""
        fixture_files = list(FIXTURES_DIR.glob("*.pptx"))
        if len(fixture_files) < 2:
            pytest.skip("Need at least 2 PPTX fixtures")
        
        for pptx_path in fixture_files[:3]:  # Test up to 3 files
            prs = Presentation(str(pptx_path))
            
            for slide in prs.slides:
                success, warnings = process_slide_safe(slide, default_config)
                assert success is True
