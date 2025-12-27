"""Pytest configuration and shared fixtures."""

import io

import pytest
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

from pp.models.config import InversionConfig


@pytest.fixture
def sample_presentation() -> Presentation:
    """Create a sample presentation with text and images for testing."""
    prs = Presentation()

    # Slide 1: Text only
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add a text box
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(1)
    
    textbox = slide1.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = "Hello World"
    
    # Add another text box with multiple paragraphs
    textbox2 = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf2 = textbox2.text_frame
    tf2.text = "First paragraph"
    p2 = tf2.add_paragraph()
    p2.text = "Second paragraph"

    # Slide 2: With image
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Create a simple test image
    img = Image.new("RGB", (100, 100), color=(255, 0, 0))  # Red image
    img_bytes = io.BytesIO()
    img.save(img_bytes, format="PNG")
    img_bytes.seek(0)
    
    slide2.shapes.add_picture(img_bytes, Inches(1), Inches(1), Inches(2), Inches(2))
    
    # Add text to slide 2 as well
    textbox3 = slide2.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    textbox3.text_frame.text = "Image slide text"

    return prs


@pytest.fixture
def sample_presentation_bytes(sample_presentation: Presentation) -> bytes:
    """Get sample presentation as bytes."""
    output = io.BytesIO()
    sample_presentation.save(output)
    output.seek(0)
    return output.read()


@pytest.fixture
def default_config() -> InversionConfig:
    """Default inversion configuration (black background, white text)."""
    return InversionConfig()


@pytest.fixture
def custom_config() -> InversionConfig:
    """Custom inversion configuration with different colors."""
    return InversionConfig(
        foreground_color=RGBColor(255, 255, 0),  # Yellow text
        background_color=RGBColor(0, 0, 128),    # Navy background
        file_suffix="(custom)",
        folder_name="Custom Output",
    )


@pytest.fixture
def sample_image() -> Image.Image:
    """Create a sample PIL Image for testing."""
    # Create a gradient image for better inversion testing
    img = Image.new("RGB", (100, 100))
    pixels = img.load()
    
    for x in range(100):
        for y in range(100):
            # Gradient from black to white
            val = int((x + y) / 2 * 255 / 100)
            pixels[x, y] = (val, val, val)  # type: ignore[index]
    
    return img


@pytest.fixture
def sample_rgba_image() -> Image.Image:
    """Create a sample RGBA Image with transparency."""
    img = Image.new("RGBA", (100, 100), (255, 0, 0, 128))  # Semi-transparent red
    return img
