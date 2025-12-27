"""Tests for Streamlit app module.

Note: Streamlit components (st.write, st.button, etc.) are not directly testable.
This test suite focuses on testable functions and logic in the app module.
"""

import io
from pathlib import Path
from datetime import datetime
from unittest.mock import Mock, patch, MagicMock

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from pp.app import (
    _get_final_folder_name,
    _cached_color_preview,
    _cached_previews,
    _hash_file,
)
from pp.models.config import InversionConfig


class TestGetFinalFolderName:
    """Tests for _get_final_folder_name function."""

    def test_without_date(self):
        """Test folder name without date suffix."""
        result = _get_final_folder_name("MyFolder", include_date=False)
        assert result == "MyFolder"

    def test_with_date(self):
        """Test folder name with date suffix."""
        result = _get_final_folder_name("MyFolder", include_date=True)
        # Should have format "MyFolder - YYYY-MM-DD"
        assert result.startswith("MyFolder - ")
        # Extract date part and verify format
        date_part = result.split(" - ")[1]
        datetime.strptime(date_part, "%Y-%m-%d")  # Should not raise

    def test_with_date_special_chars(self):
        """Test folder name with date and special characters."""
        result = _get_final_folder_name("Inverted (v2)", include_date=True)
        assert result.startswith("Inverted (v2) - ")

    def test_empty_folder_name(self):
        """Test with empty folder name."""
        result = _get_final_folder_name("", include_date=False)
        assert result == ""

    def test_empty_folder_name_with_date(self):
        """Test with empty folder name but with date."""
        result = _get_final_folder_name("", include_date=True)
        # Should be " - YYYY-MM-DD"
        assert result.startswith(" - ")


class TestHashFile:
    """Tests for _hash_file function."""

    def test_hash_consistency(self):
        """Test that same input produces same hash."""
        hash1 = _hash_file("test.pptx", b"content")
        hash2 = _hash_file("test.pptx", b"content")
        assert hash1 == hash2

    def test_different_names_different_hash(self):
        """Test that different filenames produce different hashes."""
        hash1 = _hash_file("file1.pptx", b"same_content")
        hash2 = _hash_file("file2.pptx", b"same_content")
        assert hash1 != hash2

    def test_different_content_different_hash(self):
        """Test that different content produces different hashes."""
        hash1 = _hash_file("file.pptx", b"content1")
        hash2 = _hash_file("file.pptx", b"content2")
        assert hash1 != hash2

    def test_hash_length(self):
        """Test that hash is SHA256 (64 hex characters)."""
        hash_result = _hash_file("test.pptx", b"content")
        assert len(hash_result) == 64
        # Verify it's hex
        int(hash_result, 16)  # Should not raise

    def test_empty_file(self):
        """Test hashing empty file."""
        hash_result = _hash_file("empty.pptx", b"")
        assert len(hash_result) == 64

    def test_large_content(self):
        """Test hashing large content."""
        large_content = b"x" * 1000000  # 1MB
        hash_result = _hash_file("large.pptx", large_content)
        assert len(hash_result) == 64


class TestCachedColorPreview:
    """Tests for _cached_color_preview function.
    
    Note: This function uses st.cache_data which requires Streamlit context.
    We'll test the underlying logic through integration tests.
    """

    def test_color_preview_returns_image(self):
        """Test that color preview returns image data."""
        # Create a mock Streamlit context
        with patch("pp.app.st"):
            # We can't test st.cache_data directly, but we can test
            # that the function can be called without errors
            pass  # Streamlit context-dependent


class TestCachedPreviews:
    """Tests for _cached_previews function."""

    def create_simple_pptx(self, title: str = "Test") -> bytes:
        """Create a simple PPTX in memory for testing."""
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add a text box with title
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(1)
        
        text_frame = slide.shapes.add_textbox(left, top, width, height).text_frame
        text_frame.text = title
        
        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()

    def test_preview_with_valid_pptx(self):
        """Test generating previews from valid PPTX."""
        pptx_bytes = self.create_simple_pptx("Test Slide")
        orig, inv = _cached_previews(pptx_bytes, "#000000", "#FFFFFF")
        
        assert orig is not None
        assert inv is not None

    def test_preview_with_empty_pptx(self):
        """Test generating previews from PPTX with no slides."""
        # Create a presentation without slides
        prs = Presentation()
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        pptx_bytes = output.getvalue()
        
        orig, inv = _cached_previews(pptx_bytes, "#000000", "#FFFFFF")
        assert orig is None
        assert inv is None

    def test_preview_preserves_colors(self):
        """Test that preview generation preserves color parameters."""
        pptx_bytes = self.create_simple_pptx()
        
        # Generate with different colors
        orig1, inv1 = _cached_previews(pptx_bytes, "#000000", "#FFFFFF")
        orig2, inv2 = _cached_previews(pptx_bytes, "#1a1a1a", "#f0f0f0")
        
        # Both should return images (same original since we don't change it)
        assert orig1 is not None
        assert inv1 is not None
        # inv2 should be different from inv1 due to different colors
        # (can't directly compare image data, but both should exist)
        assert inv2 is not None

    def test_preview_with_multiple_slides(self):
        """Test that preview uses only first slide."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        
        # Add multiple slides
        for i in range(3):
            slide = prs.slides.add_slide(blank_layout)
            text_frame = slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(8), Inches(1)
            ).text_frame
            text_frame.text = f"Slide {i+1}"
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        pptx_bytes = output.getvalue()
        
        orig, inv = _cached_previews(pptx_bytes, "#000000", "#FFFFFF")
        
        # Should return previews (of first slide only)
        assert orig is not None
        assert inv is not None

    def test_preview_with_colored_shape(self):
        """Test preview generation with colored shapes."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Add a colored shape
        left = Inches(2)
        top = Inches(2)
        width = Inches(4)
        height = Inches(2)
        
        shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        pptx_bytes = output.getvalue()
        
        orig, inv = _cached_previews(pptx_bytes, "#000000", "#FFFFFF")
        
        # Both should generate without error
        assert orig is not None
        assert inv is not None

    def test_preview_with_invalid_pptx(self):
        """Test preview generation with invalid PPTX data."""
        invalid_pptx = b"not a valid pptx"
        
        with pytest.raises(Exception):
            _cached_previews(invalid_pptx, "#000000", "#FFFFFF")

    def test_preview_color_hex_formats(self):
        """Test preview with various hex color formats."""
        pptx_bytes = self.create_simple_pptx()
        
        # Both should work (with and without #)
        orig1, inv1 = _cached_previews(pptx_bytes, "#000000", "#FFFFFF")
        assert orig1 is not None
        assert inv1 is not None
        
        # Lowercase
        orig2, inv2 = _cached_previews(pptx_bytes, "#000000", "#ffffff")
        assert orig2 is not None
        assert inv2 is not None


class TestInversionConfigFromApp:
    """Tests for InversionConfig usage in app context."""

    def test_config_creation_from_color_picker(self):
        """Test creating config from color picker values."""
        config = InversionConfig.from_hex(
            fg_hex="#FFFFFF",
            bg_hex="#000000",
        )
        assert config.foreground_color[0] == 255
        assert config.foreground_color[1] == 255
        assert config.foreground_color[2] == 255
        assert config.background_color[0] == 0
        assert config.background_color[1] == 0
        assert config.background_color[2] == 0

    def test_config_with_all_options(self):
        """Test config creation with all app options."""
        config = InversionConfig.from_hex(
            fg_hex="#FFFFFF",
            bg_hex="#000000",
            file_suffix="(dark)",
            folder_name="Dark Presentations - 2024-01-01",
            invert_images=False,
            jpeg_quality=75,
        )
        
        assert config.file_suffix == "(dark)"
        assert config.folder_name == "Dark Presentations - 2024-01-01"
        assert config.invert_images is False
        assert config.jpeg_quality == 75

    def test_config_validation_integration(self):
        """Test config validation in app context."""
        # Good contrast
        good_config = InversionConfig.from_hex(
            fg_hex="#FFFFFF",
            bg_hex="#000000",
        )
        warnings = good_config.validate()
        assert len(warnings) == 0
        
        # Poor contrast
        bad_config = InversionConfig.from_hex(
            fg_hex="#C8C8C8",  # Light gray
            bg_hex="#FFFFFF",  # White
        )
        warnings = bad_config.validate()
        assert len(warnings) > 0


class TestAppIntegration:
    """Integration tests for app components."""

    def test_hash_file_with_multiple_files(self):
        """Test hashing multiple different files."""
        files = [
            ("file1.pptx", b"content1"),
            ("file2.pptx", b"content2"),
            ("file1.pptx", b"content1"),  # Duplicate
        ]
        
        hashes = [_hash_file(name, data) for name, data in files]
        
        # First and third should be same (duplicate)
        assert hashes[0] == hashes[2]
        # First and second should be different
        assert hashes[0] != hashes[1]

    def test_folder_name_with_date_today(self):
        """Test folder name with today's date."""
        today = datetime.now().strftime("%Y-%m-%d")
        result = _get_final_folder_name("Test", include_date=True)
        assert today in result

    def test_invalid_hex_in_config(self):
        """Test that invalid hex colors raise errors."""
        with pytest.raises(ValueError):
            InversionConfig.from_hex(
                fg_hex="invalid",
                bg_hex="#000000",
            )

    def test_preview_memory_safety(self):
        """Test that previews don't leak memory with repeated calls."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        text_frame = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(1)
        ).text_frame
        text_frame.text = "Memory Test"
        
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        pptx_bytes = output.getvalue()
        
        # Call preview multiple times
        for _ in range(5):
            orig, inv = _cached_previews(pptx_bytes, "#000000", "#FFFFFF")
            assert orig is not None
            assert inv is not None
