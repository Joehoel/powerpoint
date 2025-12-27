"""Streamlit application for PowerPoint Inverter."""

import hashlib
import io
import logging
from datetime import datetime

import streamlit as st
from pptx import Presentation

from pp.core.inverter import process_files
from pp.models.config import InversionConfig
from pp.utils.preview import (
    generate_color_preview,
    generate_slide_preview,
    generate_slide_preview_inverted,
    hex_to_tuple,
)

logger = logging.getLogger(__name__)


def _get_final_folder_name(folder_name: str, include_date: bool) -> str:
    """Get final folder name with optional date suffix.
    
    Args:
        folder_name: Base folder name.
        include_date: Whether to append today's date.
    
    Returns:
        Folder name, potentially with date suffix (YYYY-MM-DD).
    """
    if not include_date:
        return folder_name
    date_str = datetime.now().strftime("%Y-%m-%d")
    return f"{folder_name} - {date_str}"


@st.cache_data(show_spinner=False)
def _cached_color_preview(bg_hex: str, fg_hex: str):
    return generate_color_preview(
        background_color=hex_to_tuple(bg_hex),
        foreground_color=hex_to_tuple(fg_hex),
    )


def _cached_previews(pptx_bytes: bytes, bg_hex: str, fg_hex: str):
    prs = Presentation(io.BytesIO(pptx_bytes))
    if not prs.slides:
        return None, None
    orig = generate_slide_preview(
        prs.slides[0],
        background_color=(255, 255, 255),
        text_color=(0, 0, 0),
    )
    inv = generate_slide_preview_inverted(
        prs.slides[0],
        background_color=hex_to_tuple(bg_hex),
        foreground_color=hex_to_tuple(fg_hex),
    )
    return orig, inv


def _hash_file(name: str, data: bytes) -> str:
    h = hashlib.sha256()
    h.update(name.encode())
    h.update(data)
    return h.hexdigest()


def main():

    """Main Streamlit application entry point."""
    st.set_page_config(
        page_title="PowerPoint Inverter",
        page_icon="🔄",
        layout="wide",
    )

    st.title("PowerPoint Inverter")
    st.write(
        "Invert colors in your PowerPoint presentations. "
        "Upload .pptx files or a .zip archive containing presentations."
    )

    # Initialize session state
    if "processed_result" not in st.session_state:
        st.session_state.processed_result = None
    if "preview_generated" not in st.session_state:
        st.session_state.preview_generated = False

    # Sidebar configuration
    with st.sidebar:
        st.header("Color Options")

        col1, col2 = st.columns(2)
        with col1:
            bg_color = st.color_picker(
                "Background",
                value="#000000",
                help="Color for slide backgrounds",
            )
        with col2:
            fg_color = st.color_picker(
                "Text",
                value="#FFFFFF",
                help="Color for text content",
            )

        # Check color contrast and show warnings if needed
        config = InversionConfig.from_hex(
            fg_hex=fg_color,
            bg_hex=bg_color,
        )
        contrast_warnings = config.validate()
        if contrast_warnings:
            for warning in contrast_warnings:
                st.warning(warning)

        # Show color preview
        st.subheader("Color Preview")
        color_preview = _cached_color_preview(bg_color, fg_color)
        st.image(color_preview, use_container_width=True)

        st.divider()
        st.header("Output Options")

        file_suffix = st.text_input(
            "File Suffix",
            value="(inverted)",
            help="Text added to output filenames",
        )

        folder_name = st.text_input(
            "Folder Name",
            value="Inverted Presentations",
            help="Name of the output ZIP file",
        )

        include_date = st.checkbox("Include date in folder name", value=True)

        invert_images = st.checkbox(
            "Invert image colors",
            value=True,
            help="Apply color inversion to images in slides",
        )

    # Main content area
    uploaded_files = st.file_uploader(
        "Upload PowerPoint files",
        type=["pptx", "zip"],
        accept_multiple_files=True,
        help="Select .pptx files or a .zip archive containing presentations",
    )

    if uploaded_files:
        st.success(f"Uploaded {len(uploaded_files)} file(s)")

        # File preview section
        file_blobs = [(f.name, f.getvalue()) for f in uploaded_files]
        with st.expander("Uploaded Files", expanded=False):
            for name, data in file_blobs:
                file_size = len(data) / 1024  # KB
                st.write(f"- **{name}** ({file_size:.1f} KB)")

        # Preview section
        st.subheader("Preview")

        preview_cols = st.columns(2)

        # Choose first PPTX bytes once
        first_pptx_bytes = None
        for f in uploaded_files:
            if f.name.endswith(".pptx"):
                first_pptx_bytes = f.getvalue()
                break

        preview_cache = st.session_state.get("preview_cache")
        preview_key = None
        cached_orig = cached_inv = None
        if first_pptx_bytes:
            preview_key = (
                _hash_file("preview", first_pptx_bytes),
                bg_color,
                fg_color,
            )
            if preview_cache and preview_cache.get("key") == preview_key:
                cached_orig = preview_cache.get("orig")
                cached_inv = preview_cache.get("inv")

        # Compute previews once if not cached
        if first_pptx_bytes and (cached_orig is None or cached_inv is None):
            try:
                with st.spinner("Loading preview..."):
                    orig_img, inv_img = _cached_previews(first_pptx_bytes, bg_color, fg_color)
                    st.session_state.preview_cache = {
                        "key": preview_key,
                        "orig": orig_img,
                        "inv": inv_img,
                    }
                    cached_orig, cached_inv = orig_img, inv_img
            except Exception as e:
                st.warning(f"Could not generate preview: {e}")

        with preview_cols[0]:
            st.write("**Original**")
            if first_pptx_bytes:
                if cached_orig:
                    st.image(cached_orig, width=320)
                else:
                    st.info("No slides in presentation")

        with preview_cols[1]:
            st.write("**After Inversion**")
            if first_pptx_bytes:
                if cached_inv:
                    st.image(cached_inv, width=320)
                else:
                    st.info("No slides in presentation")

        st.divider()

        # Process button
        if st.button("Invert Presentations", type="primary", use_container_width=True):
            # Build configuration
            final_folder_name = _get_final_folder_name(folder_name, include_date)

            config = InversionConfig.from_hex(
                fg_hex=fg_color,
                bg_hex=bg_color,
                file_suffix=file_suffix,
                folder_name=final_folder_name,
                invert_images=invert_images,
            )

            # Build cache key for manual reuse
            file_hashes = tuple((name, _hash_file(name, data)) for name, data in file_blobs)
            cache_key = (fg_color, bg_color, file_suffix, final_folder_name, invert_images, file_hashes)

            progress_bar = st.progress(0)
            status_text = st.empty()
            warnings_container = st.container()

            cached_entry = st.session_state.get("last_result")
            if not isinstance(cached_entry, dict):
                cached_entry = None
            cached_hit = (
                cached_entry is not None
                and cached_entry.get("key") == cache_key
                and cached_entry.get("result") is not None
            )
            if cached_hit and cached_entry is not None:
                result = cached_entry["result"]
                progress_bar.progress(1.0)
                status_text.text(
                    f"Complete! (cached) Processed {result.successful_files}/{result.total_files} files"
                )
            else:
                def update_progress(current: int, total: int, filename: str):
                    progress = current / total if total > 0 else 1.0
                    progress_bar.progress(progress)
                    status_text.text(f"Processing: {filename} ({current}/{total})")

                with st.spinner("Processing presentations..."):
                    # Reset file pointers before processing
                    for f in uploaded_files:
                        f.seek(0)
                    result = process_files(
                        uploaded_files,
                        config,
                        progress_callback=update_progress,
                    )

                progress_bar.progress(1.0)
                status_text.text(
                    f"Complete! Processed {result.successful_files}/{result.total_files} files"
                )

                st.session_state.last_result = {"key": cache_key, "result": result}

            # Store result in session state
            st.session_state.processed_result = result
            st.session_state.cached_hit = bool(cached_hit)

            # Show warnings
            if result.all_warnings:
                with warnings_container:
                    st.warning(f"Completed with {len(result.all_warnings)} warning(s)")
                    with st.expander("Show warnings"):
                        for warning in result.all_warnings:
                            st.write(f"- {warning}")

        # Show download button if we have results
        if st.session_state.processed_result is not None:
            result = st.session_state.processed_result
            cached_hit_flag = st.session_state.get("cached_hit", False)

            if result.output_zip:
                # Generate filename with date
                final_folder_name = _get_final_folder_name(folder_name, include_date)

                st.markdown("### Download")
                if cached_hit_flag:
                    st.caption("Using cached result (same files and colors)")

                st.download_button(
                    label=f"Download ({result.successful_files} files)",
                    data=result.output_zip,
                    file_name=f"{final_folder_name}.zip",
                    mime="application/zip",
                    type="primary",
                    use_container_width=True,
                )

                # Show processing summary table
                st.markdown("### Processing Summary")
                summary_rows = []
                for r in result.results:
                    summary_rows.append(
                        {
                            "File": r.filename,
                            "Status": "Success" if r.success else "Failed",
                            "Warnings": len(r.warnings),
                        }
                    )
                st.table(summary_rows)

                # Detailed warnings
                if result.all_warnings:
                    with st.expander("Warnings (detailed)"):
                        for warning in result.all_warnings:
                            st.write(f"- {warning}")
            else:
                st.error("No files were successfully processed")

    else:
        # Show instructions when no files uploaded
        st.info(
            "👆 Upload PowerPoint files to get started. "
            "You can upload multiple .pptx files or a .zip archive."
        )

        # Feature highlights
        st.subheader("Features")
        cols = st.columns(3)

        with cols[0]:
            st.markdown("**Custom Colors**")
            st.write("Choose your own foreground and background colors")

        with cols[1]:
            st.markdown("**Batch Processing**")
            st.write("Process multiple files at once with ZIP support")

        with cols[2]:
            st.markdown("**Fast & Parallel**")
            st.write("Optimized processing for large batches")


if __name__ == "__main__":
    main()
