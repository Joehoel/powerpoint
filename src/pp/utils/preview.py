"""Preview generation for PowerPoint Inverter.

This module provides utilities for generating slide previews by
rendering the actual slide content (images and text) at their
approximate positions.
"""

import io
import logging
from typing import TYPE_CHECKING

from PIL import Image, ImageDraw, ImageFont
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

if TYPE_CHECKING:
    from pptx.slide import Slide

logger = logging.getLogger(__name__)

# Default preview dimensions (4:3 aspect ratio matching typical slides)
DEFAULT_PREVIEW_WIDTH = 400
DEFAULT_PREVIEW_HEIGHT = 300


def generate_slide_preview(
    slide: "Slide",
    width: int = DEFAULT_PREVIEW_WIDTH,
    height: int = DEFAULT_PREVIEW_HEIGHT,
    background_color: tuple[int, int, int] = (255, 255, 255),
    text_color: tuple[int, int, int] = (0, 0, 0),
) -> bytes:
    """Generate a preview image of a slide by rendering its content.

    This renders actual slide content including:
    - Background color
    - Images (scaled and positioned)
    - Text (at approximate positions)

    Args:
        slide: The slide to preview.
        width: Preview image width in pixels.
        height: Preview image height in pixels.
        background_color: RGB tuple for background.
        text_color: RGB tuple for text.

    Returns:
        PNG image as bytes.
    """
    # Get slide dimensions for scaling
    presentation = slide.part.package.presentation_part.presentation
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Calculate scale factor
    scale_x = width / slide_width
    scale_y = height / slide_height
    scale = min(scale_x, scale_y)

    # Create base image with background
    img = Image.new("RGBA", (width, height), (*background_color, 255))

    # Get font for text rendering
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", max(10, int(14 * scale * 50)))
        title_font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", max(12, int(18 * scale * 50)))
    except (OSError, IOError):
        font = ImageFont.load_default()
        title_font = font

    # Collect and sort shapes by z-order (render order)
    # Pictures first, then text on top
    pictures = []
    text_shapes = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pictures.append(shape)
        elif shape.has_text_frame:
            text_shapes.append(shape)

    # Render pictures
    for shape in pictures:
        try:
            _render_picture(img, shape, scale, background_color, text_color)
        except Exception as e:
            logger.debug(f"Failed to render picture: {e}")

    # Render text on top
    draw = ImageDraw.Draw(img)
    for shape in text_shapes:
        try:
            _render_text(draw, shape, scale, text_color, font, width, height)
        except Exception as e:
            logger.debug(f"Failed to render text: {e}")

    # Add subtle border
    draw.rectangle([(0, 0), (width - 1, height - 1)], outline=(180, 180, 180))

    # Convert to bytes
    output = io.BytesIO()
    # Convert to RGB for smaller file size if no transparency needed
    img_rgb = Image.new("RGB", img.size, background_color)
    img_rgb.paste(img, mask=img.split()[3] if img.mode == "RGBA" else None)
    img_rgb.save(output, format="PNG", optimize=True)
    output.seek(0)
    return output.read()


def _render_picture(
    canvas: Image.Image,
    shape,
    scale: float,
    background_color: tuple[int, int, int],
    text_color: tuple[int, int, int],
) -> None:
    """Render a picture shape onto the canvas.

    Args:
        canvas: The PIL Image to draw on.
        shape: The picture shape from python-pptx.
        scale: Scale factor for converting EMU to pixels.
        background_color: Background color (for simulating inversion preview).
        text_color: Text/foreground color (for simulating inversion preview).
    """
    # Calculate position and size in pixels
    left = int(shape.left * scale)
    top = int(shape.top * scale)
    width = int(shape.width * scale)
    height = int(shape.height * scale)

    # Ensure minimum size
    width = max(width, 1)
    height = max(height, 1)

    # Load and resize the image
    with io.BytesIO(shape.image.blob) as img_stream:
        with Image.open(img_stream) as pic:
            # Handle transparency
            if pic.mode == "RGBA":
                pic_resized = pic.resize((width, height), Image.Resampling.LANCZOS)
            elif pic.mode == "RGB" and "transparency" in pic.info:
                # Convert transparency color to alpha
                pic_rgba = pic.convert("RGBA")
                trans_color = pic.info["transparency"]
                data = pic_rgba.getdata()
                new_data = []
                for item in data:
                    if item[:3] == trans_color:
                        new_data.append((item[0], item[1], item[2], 0))
                    else:
                        new_data.append(item)
                pic_rgba.putdata(new_data)
                pic_resized = pic_rgba.resize((width, height), Image.Resampling.LANCZOS)
            else:
                pic_rgb = pic.convert("RGBA")
                pic_resized = pic_rgb.resize((width, height), Image.Resampling.LANCZOS)

            # Paste onto canvas
            if pic_resized.mode == "RGBA":
                canvas.paste(pic_resized, (left, top), pic_resized)
            else:
                canvas.paste(pic_resized, (left, top))


def _render_text(
    draw: ImageDraw.ImageDraw,
    shape,
    scale: float,
    text_color: tuple[int, int, int],
    font: ImageFont.FreeTypeFont | ImageFont.ImageFont,
    canvas_width: int,
    canvas_height: int,
) -> None:
    """Render a text shape onto the canvas.

    Args:
        draw: PIL ImageDraw object.
        shape: The text shape from python-pptx.
        scale: Scale factor for converting EMU to pixels.
        text_color: Color for the text.
        font: Font to use for rendering.
        canvas_width: Width of the canvas.
        canvas_height: Height of the canvas.
    """
    # Get text content
    text = shape.text_frame.text.strip()  # type: ignore[union-attr]
    if not text:
        return

    # Calculate position
    left = int(shape.left * scale)
    top = int(shape.top * scale)
    width = int(shape.width * scale)

    # Clamp to canvas bounds
    left = max(0, min(left, canvas_width - 10))
    top = max(0, min(top, canvas_height - 10))

    # Simple text rendering at the shape position
    # Truncate if too long
    max_chars = max(10, int(width / 6))  # Rough estimate of chars that fit
    if len(text) > max_chars:
        text = text[:max_chars - 3] + "..."

    # Draw text with slight shadow for readability
    try:
        draw.text((left + 1, top + 1), text, fill=(128, 128, 128, 128), font=font)
        draw.text((left, top), text, fill=(*text_color, 255), font=font)
    except Exception:
        # Fallback without shadow
        draw.text((left, top), text, fill=text_color, font=font)


def generate_slide_preview_inverted(
    slide: "Slide",
    width: int = DEFAULT_PREVIEW_WIDTH,
    height: int = DEFAULT_PREVIEW_HEIGHT,
    background_color: tuple[int, int, int] = (0, 0, 0),
    foreground_color: tuple[int, int, int] = (255, 255, 255),
) -> bytes:
    """Generate a preview of how the slide will look after inversion.

    This applies color transformation to images and uses the target colors.

    Args:
        slide: The slide to preview.
        width: Preview image width in pixels.
        height: Preview image height in pixels.
        background_color: Target background color.
        foreground_color: Target foreground/text color.

    Returns:
        PNG image as bytes.
    """
    from pp.core.image_processor import apply_color_transform

    # Get slide dimensions for scaling
    presentation = slide.part.package.presentation_part.presentation
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    # Calculate scale factor
    scale_x = width / slide_width
    scale_y = height / slide_height
    scale = min(scale_x, scale_y)

    # Create base image with inverted background
    img = Image.new("RGBA", (width, height), (*background_color, 255))

    # Get font for text rendering
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", max(10, int(14 * scale * 50)))
    except (OSError, IOError):
        font = ImageFont.load_default()

    # Render shapes
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                _render_picture_inverted(
                    img, shape, scale, background_color, foreground_color, apply_color_transform
                )
            except Exception as e:
                logger.debug(f"Failed to render inverted picture: {e}")

        elif shape.has_text_frame:
            try:
                draw = ImageDraw.Draw(img)
                _render_text(draw, shape, scale, foreground_color, font, width, height)
            except Exception as e:
                logger.debug(f"Failed to render text: {e}")

    # Add subtle border
    draw = ImageDraw.Draw(img)
    draw.rectangle([(0, 0), (width - 1, height - 1)], outline=(100, 100, 100))

    # Convert to bytes
    output = io.BytesIO()
    img_rgb = Image.new("RGB", img.size, background_color)
    img_rgb.paste(img, mask=img.split()[3] if img.mode == "RGBA" else None)
    img_rgb.save(output, format="PNG", optimize=True)
    output.seek(0)
    return output.read()


def _render_picture_inverted(
    canvas: Image.Image,
    shape,
    scale: float,
    background_color: tuple[int, int, int],
    foreground_color: tuple[int, int, int],
    apply_color_transform,
) -> None:
    """Render an inverted picture shape onto the canvas.

    Args:
        canvas: The PIL Image to draw on.
        shape: The picture shape from python-pptx.
        scale: Scale factor for converting EMU to pixels.
        background_color: Target background color.
        foreground_color: Target foreground color.
        apply_color_transform: Function to apply color inversion.
    """
    # Calculate position and size in pixels
    left = int(shape.left * scale)
    top = int(shape.top * scale)
    width = int(shape.width * scale)
    height = int(shape.height * scale)

    # Ensure minimum size
    width = max(width, 1)
    height = max(height, 1)

    # Load, invert, and resize the image
    with io.BytesIO(shape.image.blob) as img_stream:
        with Image.open(img_stream) as pic:
            # Apply color inversion
            inverted = apply_color_transform(pic, background_color, foreground_color)

            # Resize
            pic_resized = inverted.resize((width, height), Image.Resampling.LANCZOS)

            # Paste onto canvas
            if pic_resized.mode == "RGBA":
                canvas.paste(pic_resized, (left, top), pic_resized)
            else:
                canvas.paste(pic_resized, (left, top))


def generate_color_preview(
    background_color: tuple[int, int, int],
    foreground_color: tuple[int, int, int],
    width: int = 200,
    height: int = 100,
) -> bytes:
    """Generate a preview showing the color scheme.

    Creates an image with the background color and sample text
    in the foreground color.

    Args:
        background_color: RGB tuple for background.
        foreground_color: RGB tuple for text.
        width: Preview width in pixels.
        height: Preview height in pixels.

    Returns:
        PNG image as bytes.
    """
    img = Image.new("RGB", (width, height), background_color)
    draw = ImageDraw.Draw(img)

    # Try to get a font
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 16)
    except (OSError, IOError):
        font = ImageFont.load_default()

    # Draw sample text
    sample_text = "Sample Text"
    bbox = draw.textbbox((0, 0), sample_text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]

    x = (width - text_width) // 2
    y = (height - text_height) // 2

    draw.text((x, y), sample_text, fill=foreground_color, font=font)

    # Add border
    draw.rectangle([(0, 0), (width - 1, height - 1)], outline=(128, 128, 128))

    output = io.BytesIO()
    img.save(output, format="PNG")
    output.seek(0)
    return output.read()


def rgb_color_to_tuple(color: RGBColor) -> tuple[int, int, int]:
    """Convert RGBColor to tuple.

    Args:
        color: pptx RGBColor instance.

    Returns:
        Tuple of (red, green, blue) integers.
    """
    return (color[0], color[1], color[2])


def hex_to_tuple(hex_color: str) -> tuple[int, int, int]:
    """Convert hex color string to RGB tuple.

    Args:
        hex_color: Color as hex string (e.g., "#FF0000").

    Returns:
        Tuple of (red, green, blue) integers.
    """
    hex_color = hex_color.lstrip("#")
    return (
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )
