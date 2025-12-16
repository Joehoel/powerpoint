"""Image processing for PowerPoint Inverter.

This module handles memory-efficient image color transformation.
"""

import io
import logging
from typing import TYPE_CHECKING

from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

if TYPE_CHECKING:
    from pptx.shapes.picture import Picture
    from pptx.slide import Slide

logger = logging.getLogger(__name__)


def apply_color_transform(
    img: Image.Image,
    target_dark: tuple[int, int, int],
    target_light: tuple[int, int, int],
) -> Image.Image:
    """Apply color transformation to map dark colors to target_dark and light to target_light.

    This creates a more sophisticated inversion that maps the image's color range
    to the specified target colors rather than simple RGB inversion.

    Handles three types of transparency:
    1. RGBA images with explicit alpha channel
    2. RGB images with transparency color defined in img.info['transparency']
    3. Palette images with transparency

    Args:
        img: PIL Image to transform.
        target_dark: RGB tuple for dark colors (e.g., background color).
        target_light: RGB tuple for light colors (e.g., foreground color).

    Returns:
        Transformed PIL Image (RGBA if transparency present, RGB otherwise).
    """
    from PIL import ImageOps

    alpha = None
    transparency_color = None

    # Handle different image modes and transparency types
    if img.mode == "RGBA":
        # Explicit alpha channel
        alpha = img.split()[3]
        img_rgb = img.convert("RGB")
    elif img.mode == "RGB" and "transparency" in img.info:
        # RGB with transparency color - convert to RGBA first
        transparency_color = img.info["transparency"]
        img_rgb, alpha = _convert_transparency_color_to_alpha(img, transparency_color)
    elif img.mode == "P":
        # Palette mode - convert to RGBA to handle any transparency
        rgba = img.convert("RGBA")
        alpha = rgba.split()[3]
        img_rgb = rgba.convert("RGB")
    elif img.mode != "RGB":
        img_rgb = img.convert("RGB")
    else:
        img_rgb = img

    # Invert the image
    inverted = ImageOps.invert(img_rgb)

    # If target colors are standard black/white inversion, we're done
    if target_dark == (0, 0, 0) and target_light == (255, 255, 255):
        result = inverted
    else:
        # Apply color mapping: remap inverted image to target color range
        result = _remap_colors(inverted, target_dark, target_light)

    # Restore alpha channel if present
    if alpha is not None:
        result = result.convert("RGBA")
        result.putalpha(alpha)

    return result


def _convert_transparency_color_to_alpha(
    img: Image.Image,
    transparency_color: tuple[int, int, int],
) -> tuple[Image.Image, Image.Image]:
    """Convert an RGB image with transparency color to RGB + alpha channel.

    PNG images can define transparency via a specific color rather than an alpha
    channel. This function creates an explicit alpha channel where pixels matching
    the transparency color become fully transparent.

    Args:
        img: RGB PIL Image with transparency color.
        transparency_color: RGB tuple that represents transparent pixels.

    Returns:
        Tuple of (RGB image, alpha channel as grayscale Image).
    """
    import numpy as np

    # Convert to numpy for efficient comparison
    arr = np.array(img)

    # Create alpha channel: 255 for opaque, 0 for transparent
    # Pixels matching transparency color become transparent
    matches = np.all(arr == transparency_color, axis=2)
    alpha_arr = np.where(matches, 0, 255).astype(np.uint8)

    alpha = Image.fromarray(alpha_arr, mode="L")
    return img, alpha


def _remap_colors(
    img: Image.Image,
    target_dark: tuple[int, int, int],
    target_light: tuple[int, int, int],
) -> Image.Image:
    """Remap image colors from black-white range to target color range.

    Args:
        img: RGB PIL Image (inverted).
        target_dark: Target color for dark values.
        target_light: Target color for light values.

    Returns:
        Remapped PIL Image.
    """
    import numpy as np

    # Convert to numpy array for efficient processing
    arr = np.array(img, dtype=np.float32)

    # Normalize to 0-1 range
    arr = arr / 255.0

    # For each channel, interpolate between target_dark and target_light
    result = np.zeros_like(arr)
    for i in range(3):
        dark_val = target_dark[i] / 255.0
        light_val = target_light[i] / 255.0
        result[:, :, i] = dark_val + arr[:, :, i] * (light_val - dark_val)

    # Convert back to 0-255 range
    result = (result * 255).astype(np.uint8)
    return Image.fromarray(result, mode="RGB")


def invert_image(
    slide: "Slide",
    shape: "Picture",
    background_color: RGBColor,
    foreground_color: RGBColor,
) -> str | None:
    """Invert colors of a picture shape on a slide.

    This function extracts the image, applies color transformation,
    and replaces the original with the transformed version.

    Args:
        slide: The slide containing the picture.
        shape: The picture shape to invert.
        background_color: Target color for originally light areas.
        foreground_color: Target color for originally dark areas.

    Returns:
        Warning message if processing failed, None on success.
    """
    try:
        # Extract image data
        image_blob = shape.image.blob

        # Process image with context manager for memory efficiency
        with io.BytesIO(image_blob) as image_stream:
            with Image.open(image_stream) as img:
                # Apply color transformation
                # Note: For inversion, light areas become dark (background)
                # and dark areas become light (foreground)
                transformed = apply_color_transform(
                    img,
                    target_dark=(background_color[0], background_color[1], background_color[2]),
                    target_light=(foreground_color[0], foreground_color[1], foreground_color[2]),
                )

                # Save to bytes
                output_stream = io.BytesIO()
                # Determine format - use PNG for quality, or original format if possible
                save_format = "PNG"
                transformed.save(output_stream, format=save_format)
                output_stream.seek(0)

                # Store position and size before removing
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Remove original shape
                shape_element = shape._element
                shape_element.getparent().remove(shape_element)

                # Add new image at same position
                slide.shapes.add_picture(output_stream, left, top, width, height)

        return None

    except Exception as e:
        logger.warning(f"Failed to process image: {e}")
        return f"Image processing failed: {e}"


def is_picture_shape(shape) -> bool:
    """Check if a shape is a picture that can be inverted.

    Args:
        shape: A pptx shape object.

    Returns:
        True if shape is a processable picture.
    """
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
