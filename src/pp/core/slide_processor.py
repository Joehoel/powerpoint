"""Slide processing for PowerPoint Inverter.

This module handles single-pass slide processing for optimal performance.
"""

import logging
from typing import TYPE_CHECKING

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pp.core.image_processor import invert_image, is_picture_shape
from pp.models.config import InversionConfig

if TYPE_CHECKING:
    from pptx.slide import Slide

logger = logging.getLogger(__name__)


def process_slide(slide: "Slide", config: InversionConfig) -> list[str]:
    """Process a single slide, inverting colors in one pass.

    This function performs all color inversions (background, text, images)
    in a single iteration through the slide's shapes for optimal performance.

    Args:
        slide: The slide to process.
        config: Inversion configuration with target colors.

    Returns:
        List of warning messages for any issues encountered.
    """
    warnings: list[str] = []

    # Set slide background color
    try:
        _set_background_color(slide, config.background_color)
    except Exception as e:
        logger.warning(f"Failed to set background: {e}")
        warnings.append(f"Background: {e}")

    # Collect picture shapes first (we'll modify the shapes collection)
    picture_shapes = []

    # Single pass through all shapes
    for shape in slide.shapes:
        # Process text in shapes
        if shape.has_text_frame:
            try:
                _invert_text_color(shape, config.foreground_color)
            except Exception as e:
                logger.warning(f"Failed to invert text in shape: {e}")
                warnings.append(f"Text: {e}")

        # Collect pictures for later processing
        # (Can't modify during iteration)
        if is_picture_shape(shape):
            picture_shapes.append(shape)

    # Process collected pictures
    if config.invert_images:
        for shape in picture_shapes:
            warning = invert_image(
                slide, shape, config.background_color, config.foreground_color
            )
            if warning:
                warnings.append(warning)

    return warnings


def _set_background_color(slide: "Slide", color: RGBColor) -> None:
    """Set the slide background to a solid color.

    Args:
        slide: The slide to modify.
        color: The background color.
    """
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _invert_text_color(shape, color: RGBColor) -> None:
    """Set all text in a shape to the specified color.

    Args:
        shape: A shape with a text frame.
        color: The text color to apply.
    """
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = color


def process_slide_safe(slide: "Slide", config: InversionConfig) -> tuple[bool, list[str]]:
    """Process a slide with full error handling.

    Args:
        slide: The slide to process.
        config: Inversion configuration.

    Returns:
        Tuple of (success: bool, warnings: list[str]).
    """
    try:
        warnings = process_slide(slide, config)
        return True, warnings
    except Exception as e:
        logger.exception(f"Slide processing failed: {e}")
        return False, [f"Slide failed: {e}"]
