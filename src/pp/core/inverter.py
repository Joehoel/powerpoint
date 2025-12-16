"""Main orchestration for PowerPoint Inverter.

This module handles processing of presentations and files.
Sequential processing is used to minimize memory usage on resource-constrained
environments like small VPS instances.
"""

import io
import logging
import os
import zipfile
from concurrent.futures import FIRST_COMPLETED, ProcessPoolExecutor, ThreadPoolExecutor, wait
from dataclasses import dataclass
from pathlib import Path
from typing import IO, TYPE_CHECKING, Callable, Iterator, cast

from pptx import Presentation
from pptx.dml.color import RGBColor

from pp.core.slide_processor import process_slide_safe
from pp.models.config import InversionConfig

if TYPE_CHECKING:
    from pptx.presentation import Presentation as PresentationType
    from streamlit.runtime.uploaded_file_manager import UploadedFile

logger = logging.getLogger(__name__)

ConfigPayload = tuple[tuple[int, int, int], tuple[int, int, int], str, str, bool]


@dataclass
class ProcessingResult:
    """Result of processing a single file."""

    filename: str
    success: bool
    output_data: bytes | None  # Store bytes directly instead of path
    warnings: list[str]


@dataclass
class BatchResult:
    """Result of processing multiple files."""

    results: list[ProcessingResult]
    output_zip: bytes
    total_files: int
    successful_files: int

    @property
    def all_warnings(self) -> list[str]:
        """Get all warnings from all files."""
        warnings = []
        for result in self.results:
            if result.warnings:
                for warning in result.warnings:
                    warnings.append(f"{result.filename}: {warning}")
        return warnings


ProgressCallback = Callable[[int, int, str], None]


def _get_max_workers() -> int:
    """Return a small, safe worker count for constrained hosts."""
    cpu = os.cpu_count() or 1
    return max(1, min(2, cpu))


def _use_threads() -> bool:
    """Allow forcing threads via env to reduce memory pressure."""
    return os.environ.get("PP_FORCE_THREADS", "0") == "1"


def _color_to_tuple(color: RGBColor) -> tuple[int, int, int]:
    return (int(color[0]), int(color[1]), int(color[2]))


def _config_to_payload(config: InversionConfig) -> ConfigPayload:
    return (
        _color_to_tuple(config.foreground_color),
        _color_to_tuple(config.background_color),
        config.file_suffix,
        config.folder_name,
        config.invert_images,
    )


def _config_from_payload(
    payload: tuple[tuple[int, int, int], tuple[int, int, int], str, str, bool]
) -> InversionConfig:
    fg, bg, suffix, folder, invert_images = payload
    return InversionConfig(
        foreground_color=RGBColor(*fg),
        background_color=RGBColor(*bg),
        file_suffix=suffix,
        folder_name=folder,
        invert_images=invert_images,
    )


def _process_file(
    file_data: bytes,
    filename: str,
    config: InversionConfig | tuple[tuple[int, int, int], tuple[int, int, int], str, str, bool],
) -> ProcessingResult:
    """Process a single PPTX file.

    Args:
        file_data: Raw bytes of the PPTX file.
        filename: Original filename.
        config: Inversion configuration.

    Returns:
        ProcessingResult with output bytes.
    """
    cfg = config if isinstance(config, InversionConfig) else _config_from_payload(config)

    try:
        # Load presentation from bytes
        with io.BytesIO(file_data) as file_stream:
            prs = Presentation(file_stream)

            # Process all slides
            all_warnings: list[str] = []
            for idx, slide in enumerate(prs.slides):
                success, warnings = process_slide_safe(slide, cfg)
                if warnings:
                    for warning in warnings:
                        all_warnings.append(f"Slide {idx + 1}: {warning}")

            # Generate output filename
            name_without_ext = os.path.splitext(filename)[0]
            output_filename = f"{name_without_ext} {cfg.file_suffix}.pptx"

            # Save to bytes
            output_stream = io.BytesIO()
            prs.save(output_stream)
            output_data = output_stream.getvalue()

            return ProcessingResult(
                filename=output_filename,
                success=True,
                output_data=output_data,
                warnings=all_warnings,
            )

    except Exception as e:
        logger.exception(f"Failed to process {filename}: {e}")
        return ProcessingResult(
            filename=filename,
            success=False,
            output_data=None,
            warnings=[f"Processing failed: {e}"],
        )


def process_presentation(
    prs: "PresentationType",
    config: InversionConfig,
    progress_callback: ProgressCallback | None = None,
    filename: str = "presentation",
) -> list[str]:
    """Process all slides in a presentation.

    Args:
        prs: The Presentation object to process.
        config: Inversion configuration.
        progress_callback: Optional callback(current, total, status) for progress.
        filename: Name of the file being processed (for logging).

    Returns:
        List of warning messages.
    """
    all_warnings: list[str] = []
    total_slides = len(prs.slides)

    if total_slides == 0:
        return ["Presentation has no slides"]

    for idx, slide in enumerate(prs.slides):
        success, warnings = process_slide_safe(slide, config)

        if warnings:
            for warning in warnings:
                all_warnings.append(f"Slide {idx + 1}: {warning}")

        if progress_callback:
            progress_callback(idx + 1, total_slides, f"Processing slide {idx + 1}/{total_slides}")

    return all_warnings


def process_single_file(
    file_data: bytes,
    filename: str,
    config: InversionConfig,
    output_dir: Path,
    progress_callback: ProgressCallback | None = None,
) -> ProcessingResult:
    """Process a single PPTX file and save to disk.

    Args:
        file_data: Raw bytes of the PPTX file.
        filename: Original filename.
        config: Inversion configuration.
        output_dir: Directory to save output file.
        progress_callback: Optional progress callback.

    Returns:
        ProcessingResult with success status and any warnings.
    """
    try:
        # Load presentation from bytes
        with io.BytesIO(file_data) as file_stream:
            prs = Presentation(file_stream)

            # Process the presentation
            warnings = process_presentation(prs, config, progress_callback, filename)

            # Generate output filename
            name_without_ext = os.path.splitext(filename)[0]
            output_filename = f"{name_without_ext} {config.file_suffix}.pptx"
            output_path = output_dir / output_filename

            # Save the modified presentation
            prs.save(str(output_path))

            return ProcessingResult(
                filename=filename,
                success=True,
                output_data=None,
                warnings=warnings,
            )

    except Exception as e:
        logger.exception(f"Failed to process {filename}: {e}")
        return ProcessingResult(
            filename=filename,
            success=False,
            output_data=None,
            warnings=[f"Processing failed: {e}"],
        )


def process_files_streaming(
    uploaded_files: list["UploadedFile"],
    config: InversionConfig,
    progress_callback: Callable[[int, int, str], None] | None = None,
) -> Iterator[ProcessingResult]:
    """Process multiple uploaded files with streaming results.

    Uses a small, bounded executor (processes by default, threads if
    PP_FORCE_THREADS=1) so at most a couple of files are in-flight.
    """
    # Collect all PPTX files to process
    files_to_process: list[tuple[bytes, str]] = []

    for uploaded_file in uploaded_files:
        if uploaded_file.name.endswith(".pptx"):
            files_to_process.append((uploaded_file.read(), uploaded_file.name))
            uploaded_file.seek(0)

        elif uploaded_file.name.endswith(".zip"):
            extracted = _extract_pptx_from_zip(uploaded_file)
            files_to_process.extend(extracted)

    total_files = len(files_to_process)

    if total_files == 0:
        return

    max_workers = _get_max_workers()
    executor_cls = ThreadPoolExecutor if _use_threads() else ProcessPoolExecutor
    logger.info(
        f"Processing {total_files} files with {executor_cls.__name__} ({max_workers} workers)"
    )

    iterator = iter(files_to_process)
    payload = _config_to_payload(config)

    def submit_next(executor, futures):
        try:
            file_data, filename = next(iterator)
        except StopIteration:
            return False
        fut = executor.submit(_process_file, file_data, filename, payload)
        futures[fut] = filename
        return True

    with executor_cls(max_workers=max_workers) as executor:
        futures: dict = {}
        for _ in range(max_workers):
            if not submit_next(executor, futures):
                break

        completed = 0
        while futures:
            done, _ = wait(futures.keys(), return_when=FIRST_COMPLETED)
            for fut in done:
                filename = futures.pop(fut)
                try:
                    result = fut.result()
                except Exception as e:
                    logger.exception(f"Future failed for {filename}: {e}")
                    result = ProcessingResult(
                        filename=filename,
                        success=False,
                        output_data=None,
                        warnings=[f"Unexpected error: {e}"],
                    )

                yield result

                completed += 1
                if progress_callback:
                    progress_callback(completed, total_files, filename)

                submit_next(executor, futures)


def process_files(
    uploaded_files: list["UploadedFile"],
    config: InversionConfig,
    progress_callback: Callable[[int, int, str], None] | None = None,
) -> BatchResult:
    """Process multiple uploaded files (PPTX or ZIP).

    Uses sequential processing to minimize memory usage on resource-constrained
    environments. The image processing optimizations (JPEG encoding, optimized
    numpy operations) provide the main performance improvements.

    Args:
        uploaded_files: List of Streamlit UploadedFile objects.
        config: Inversion configuration.
        progress_callback: Optional callback(current_file, total_files, filename).

    Returns:
        BatchResult with output ZIP and processing statistics.
    """
    # Collect all PPTX files to process
    files_to_process: list[tuple[bytes, str]] = []

    for uploaded_file in uploaded_files:
        if uploaded_file.name.endswith(".pptx"):
            files_to_process.append((uploaded_file.read(), uploaded_file.name))
            uploaded_file.seek(0)

        elif uploaded_file.name.endswith(".zip"):
            extracted = _extract_pptx_from_zip(uploaded_file)
            files_to_process.extend(extracted)

    total_files = len(files_to_process)

    if total_files == 0:
        return BatchResult(
            results=[],
            output_zip=b"",
            total_files=0,
            successful_files=0,
        )

    results: list[ProcessingResult] = []
    executor_cls = ThreadPoolExecutor if _use_threads() else ProcessPoolExecutor
    max_workers = _get_max_workers()
    logger.info(
        f"Processing {total_files} files with {executor_cls.__name__} ({max_workers} workers)"
    )

    iterator = iter(files_to_process)
    payload = _config_to_payload(config)

    def submit_next(executor, futures):
        try:
            file_data, filename = next(iterator)
        except StopIteration:
            return False
        fut = executor.submit(_process_file, file_data, filename, payload)
        futures[fut] = filename
        return True

    with executor_cls(max_workers=max_workers) as executor:
        futures: dict = {}
        for _ in range(max_workers):
            if not submit_next(executor, futures):
                break

        completed = 0
        while futures:
            done, _ = wait(futures.keys(), return_when=FIRST_COMPLETED)
            for fut in done:
                filename = futures.pop(fut)
                try:
                    result = fut.result()
                except Exception as e:
                    logger.exception(f"Future failed for {filename}: {e}")
                    result = ProcessingResult(
                        filename=filename,
                        success=False,
                        output_data=None,
                        warnings=[f"Unexpected error: {e}"],
                    )
                results.append(result)

                completed += 1
                if progress_callback:
                    progress_callback(completed, total_files, filename)

                submit_next(executor, futures)

    # Create output ZIP from in-memory results
    output_zip = _create_output_zip_from_results(results)

    successful = sum(1 for r in results if r.success)

    return BatchResult(
        results=results,
        output_zip=output_zip,
        total_files=total_files,
        successful_files=successful,
    )


def _extract_pptx_from_zip(zip_file: IO[bytes]) -> list[tuple[bytes, str]]:
    """Extract PPTX files from a ZIP archive.

    Args:
        zip_file: File-like object containing ZIP data.

    Returns:
        List of (file_data, filename) tuples.
    """
    files = []
    try:
        with zipfile.ZipFile(zip_file, "r") as zf:
            for name in zf.namelist():
                if name.endswith(".pptx") and not name.startswith("__MACOSX"):
                    data = zf.read(name)
                    # Use just the filename, not the full path
                    filename = os.path.basename(name)
                    if filename:  # Skip if it was a directory
                        files.append((data, filename))
    except zipfile.BadZipFile as e:
        logger.warning(f"Invalid ZIP file: {e}")
    return files


def _create_output_zip_from_results(results: list[ProcessingResult]) -> bytes:
    """Create a ZIP file from in-memory processing results.

    Args:
        results: List of processing results with output_data bytes.

    Returns:
        ZIP file as bytes.
    """
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as zf:
        for result in results:
            if result.success and result.output_data:
                zf.writestr(result.filename, result.output_data)

    output.seek(0)
    return output.read()
